from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import copy
import os
import inputs
DIR_PATH = os.path.dirname(os.path.realpath(__file__))
import io
import json
from PIL import Image
from collections import Counter
from pptx.dml.color import RGBColor
from colorthief import ColorThief


# Import data from the index.json file
with open('index.json') as f:
    data = json.load(f)

def CopySlide(copyFromPres, slideIndex, pasteIntoPres):
    # Specify the slide you want to copy the contents from
    slide_to_copy = copyFromPres.slides[slideIndex]

    # Define the layout you want to use for your generated pptx
    slide_layout = pasteIntoPres.slide_layouts.get_by_name("Blank")  # Ensure blank layout
    new_slide = pasteIntoPres.slides.add_slide(slide_layout)

    # Dictionary to store image and shape data
    imgDict = {}
    shapeDict = {}

    # First pass: Identify all shapes and classify them into images and shapes
    for shp in slide_to_copy.shapes:
        if shp.shape_type == 13:  # If it's an image
            try:
                img_path = os.path.join(os.getcwd(), shp.name + '.jpg')
                with open(img_path, 'wb') as f:
                    f.write(shp.image.blob)

                # Store image information for adding to the new slide
                imgDict[img_path] = [shp.left, shp.top, shp.width, shp.height]

            except Exception as e:
                print(f"Error processing image {shp.name}: {e}")
        else:
            # Classify the rest as shapes (possibly including the grey rectangle or black square)
            shapeDict[shp.shape_id] = (shp, [shp.left, shp.top, shp.width, shp.height])

    # Sort all shapes and images by their area (width * height), largest first
    all_elements = []
    for img_path, dims in imgDict.items():
        all_elements.append(('image', img_path, dims))
    for shape_id, (shape, dims) in shapeDict.items():
        all_elements.append(('shape', shape, dims))
    
    # Sort by area (width * height), largest first
    all_elements.sort(key=lambda x: x[2][2] * x[2][3], reverse=True)

    # Add the shapes and images to the new slide, starting with the largest
    for element_type, element, dims in all_elements:
        if element_type == 'image':
            # Add images to the slide
            try:
                new_slide.shapes.add_picture(element, dims[0], dims[1], dims[2], dims[3])
            except Exception as e:
                print(f"Error adding image {element} to the slide: {e}")
            finally:
                if os.path.exists(element):
                    os.remove(element)  # Remove image file after it's added to the slide
        elif element_type == 'shape':
            # Add shapes to the slide
            try:
                el = element.element
                newel = copy.deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
            except Exception as e:
                print(f"Error processing shape {element.name}: {e}")

    return new_slide


def NewSlide(index, link, outputPres):
    templatePres = Presentation(link)
    pastedSlide = CopySlide(templatePres,index,outputPres)

    #Modif slide by slide
    if index == 1 and link == "Slides_DEBUT_FIN.pptx":
        first_image = True
        for shape in pastedSlide.shapes:
            if shape.shape_type == 13 and first_image == True: #image
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                shape._element.getparent().remove(shape._element)
                pastedSlide.shapes.add_picture(inputs.input_0_A, left, top, width=width, height=height)
                first_image = False
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "DATE" in run.text:
                            # Replace the text, but keep the formatting
                            run.text = "{}".format(inputs.input_0_B)

    if index == 2 and link == "Slides_DEBUT_FIN.pptx":
        for shape in pastedSlide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "ENTREPRISE" in run.text:
                            # Replace the text, but keep the formatting
                            run.text = "{} souhaite organiser un séminaire qui réunira {} collaborateurs pour un moment de partage et d’échanges, sur {} jours et {} nuits.".format(inputs.input_1_A, inputs.input_1_B, inputs.input_1_C, inputs.input_1_D)
        for shape in pastedSlide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "DATE" in run.text:
                            # Replace the text, but keep the formatting
                            run.text = "{}".format(inputs.input_1_E)

    if index == 5 and link == "liste_HOTELS.pptx":
        for shape in pastedSlide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "NOMBRE CHAMBRES" in run.text:
                            # Replace the text, but keep the formatting
                            run.text = "Vous aurez à votre disposition {}".format(inputs.input_koutoubia_5)


    if index == 13 and link == "liste_HOTELS.pptx":
        for shape in pastedSlide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "NOMBRE CHAMBRES" in run.text:
                            # Replace the text, but keep the formatting
                            run.text = "Vous aurez à votre disposition {}".format(inputs.input_sofitel_3)

    if index == 10 and link == "Slides_DEBUT_FIN.pptx":
        for shape in pastedSlide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "CLIENT" in run.text:
                            # Replace the text, but keep the formatting
                            run.text = "{}".format(inputs.input_1_A)




def select_image_from_index(activity_label):
    """
    Selects an image from the reference PowerPoint based on the activity_label.
    activity_label: The label of the activity that maps to a specific image.
    Returns the extracted image blob.
    """
    if activity_label not in data:
        raise ValueError(f"Activity '{activity_label}' not found in the image index map")

    slide_number = data[activity_label]["image"][0]  # Get slide number from the index
    image_index = data[activity_label]["image"][1]  # Get the image index for that activity

    ref_presentation = Presentation("liste_JOUR_N.pptx")
    slide_to_copy = ref_presentation.slides[slide_number - 1]  # Convert to zero-based index
    
    picture_shapes = []
    
    # Iterate through the shapes on the slide to find pictures
    for shape in slide_to_copy.shapes:
        if shape.shape_type == 13:  # Shape type 13 means it's an image
            picture_shapes.append(shape)

    # Ensure the picture index exists
    if image_index - 1 >= len(picture_shapes):
        raise ValueError(f"Image index {image_index} is out of bounds for slide {slide_number}")

    selected_image = picture_shapes[image_index - 1]  # Select the appropriate image by index

    # Return the image blob (binary content of the image)
    return selected_image.image.blob





def CopyAndModifySlide(slideIndex, pasteIntoPres, etapes, day_index, day_date):
    """
    Copies a slide from a presentation and modifies it by replacing the existing images
    with images corresponding to the given 'etapes' (activities), maintaining the original layout.
    etapes: List of activity labels.
    """
    copyFromPres = Presentation("liste_JOUR_N.pptx")

    # Specify the slide you want to copy the contents from
    slide_to_copy = copyFromPres.slides[slideIndex]

    # Define the layout you want to use for your generated pptx
    slide_layout = pasteIntoPres.slide_layouts.get_by_name("Blank")  # Ensure blank layout
    new_slide = pasteIntoPres.slides.add_slide(slide_layout)

    # Step 1: Copy everything from the original slide (text, images, logos, etc.)
    for shp in slide_to_copy.shapes:
        el = shp.element
        newel = copy.deepcopy(el)  # Copy each element in the slide
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    # Step 2: Save the two highest-positioned images (logos) for reinsertion later
    image_shapes = []
    logo_shapes = []

    # Identify images and sort them based on their "top" position (Y-axis)
    for shp in slide_to_copy.shapes:
        if shp.shape_type == 13:  # If it's an image
            image_shapes.append(shp)

    # Sort the images by their top position (Y-axis) in ascending order
    image_shapes.sort(key=lambda x: x.top)

    # Save the two images that are highest on the slide (logos)
    for i in range(2):
        logo_shapes.append({
            "blob": image_shapes[i].image.blob,
            "left": image_shapes[i].left,
            "top": image_shapes[i].top,
            "width": image_shapes[i].width,
            "height": image_shapes[i].height
        })

    # Step 1: Extract the first two images (logos) and remove them from the image_shapes list
    logo_shapes = image_shapes[:2]  # First two images are logos
    non_logo_shapes = image_shapes[2:]  # Remaining images to work with

    # Step 2: Sort the non-logo images based on the 'left' value (from left to right)
    non_logo_shapes.sort(key=lambda shp: shp.left)

    # Step 3: Replace the images for the etapes starting from the 3rd image (excluding logos)
    for i, etape in enumerate(etapes):
        if i < len(non_logo_shapes):  # Ensure there's an image to replace
            try:
                # Get the corresponding image for this etape from the reference PowerPoint
                image_blob = select_image_from_index(etape)

                # Save the image blob to a temporary file
                temp_img_path = f'temp_image_{etape}_{i}.jpg'
                with open(temp_img_path, 'wb') as temp_img_file:
                    temp_img_file.write(image_blob)

                # Get the position and size of the current image in the slide (after sorting)
                left = non_logo_shapes[i].left
                top = non_logo_shapes[i].top
                width = non_logo_shapes[i].width
                height = non_logo_shapes[i].height


                # Remove the existing image
                non_logo_shapes[i]._element.getparent().remove(non_logo_shapes[i]._element)

                # Add the new image to the slide, using the same position and size as the original
                new_slide.shapes.add_picture(temp_img_path, left, top, width=width, height=height)

                # Clean up the temporary image file
                if os.path.exists(temp_img_path):
                    os.remove(temp_img_path)

            except Exception as e:
                print(f"Error adding image for etape '{etape}': {e}")


    # Step 4: Replace the text placeholders (ACTION 1, ACTION 2, etc.) with the etape names
    text_shapes = []
    for shp in new_slide.shapes:
        if shp.has_text_frame:
            for paragraph in shp.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_shapes.append(shp)

    for i, etape in enumerate(etapes):
        if "DINER A L'HOTEL" in etape:
            etape = "DINER A L'HOTEL"
        if i < len(text_shapes):
            for text_shape in text_shapes:
                for paragraph in text_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        action = f"ACTION {i+1}"
                        if action in run.text:  # Ensure we are replacing the right placeholders
                            run.text = etape  # Replace the placeholder text with the etape
                        if "JOUR N" in run.text:
                            run.text = f"JOUR {day_index} : "
                        if "DATE" in run.text:
                            run.text = day_date

    # Step 5: Remove any images that overlap with the logos
    for logo in logo_shapes:
        for shape in new_slide.shapes:
            if shape.shape_type == 13:  # Check if it's an image
                # Check if the image overlaps with the logo position
                if shape.left == logo.left and shape.top == logo.top:
                    shape._element.getparent().remove(shape._element)

    # Step 6: Reinsert the saved logo images at their original positions
    try:
        # Re-add the first logo (company logo) back to the slide
        temp_logo_path = f'temp_logo_{2}.jpg'
        
        # Save the logo image to a temporary file
        with open(temp_logo_path, 'wb') as temp_logo_file:
            temp_logo_file.write(logo_shapes[1].image.blob)  # Access the image blob directly from the shape

        # Add the logo image back into the slide
        new_slide.shapes.add_picture(temp_logo_path, logo_shapes[1].left, logo_shapes[1].top, logo_shapes[1].width, logo_shapes[1].height)

        # Clean up the temporary logo file
        if os.path.exists(temp_logo_path):
            os.remove(temp_logo_path)

    except Exception as e:
        print(f"Error adding logo 1: {e}")

    # Step 7: Reinsert the client logo (input logo)
    try:
        # Re-add the second logo (client logo) back to the slide
        temp_logo_path = f'temp_logo_{1}.jpg'
        
        # Add the input logo (from inputs.input_0_A) back to the slide
        new_slide.shapes.add_picture(inputs.input_0_A, logo_shapes[0].left, logo_shapes[0].top, logo_shapes[0].width, logo_shapes[0].height)

        # Clean up the temporary logo file
        if os.path.exists(temp_logo_path):
            os.remove(temp_logo_path)

    except Exception as e:
        print(f"Error adding logo 2: {e}")


    return new_slide




def get_main_color(image_path):
    """Extracts the most common RGB color in an image."""
    image = Image.open(image_path)
    image = image.convert('RGB')  # Ensure the image is converted to RGB to avoid alpha channels
    image = image.resize((100, 100))  # Resize for faster processing
    
    pixels = image.getdata()  # Get all pixel data
    most_common_color = Counter(pixels).most_common(1)[0][0]  # Get the most common color
    
    return most_common_color  # This will now return a tuple like (R, G, B)


def change_shape_colors(slide, main_color):
    """Changes the color of shapes that are not black (excluding pictures, texts, and connectors)."""
    r, g, b = main_color  # Extract RGB components from the main color

    for shape in slide.shapes:
        # Skip pictures and text, and now skip connectors (connector is shape_type 9)
        if shape.shape_type in [9, 13, 14]:  # 9 = connector, 13 = picture, 14 = text
            continue

        # Ensure the shape has a fill property and it's not None
        if hasattr(shape, 'fill') and shape.fill.type is not None:
            # Get the fill color
            fill_color = shape.fill.fore_color.rgb if shape.fill.fore_color else None

            # Check if the color is not black
            if fill_color and fill_color != RGBColor(0, 0, 0):
                # Change the shape's fill to the main color extracted from the logo
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(r, g, b)
            else:
                print(f"Skipped shape with black fill color or no color.")
        else:
            print("Skipped shape without a fill property.")




def make_agenda(pasteIntoPres, activities_by_day):
    """
    Modifies the agenda slide(s) based on the number of days and activities.

    activities_by_day: A list of lists where each inner list contains activities for that day.
                       Historically the caller passed an extra trailing element; to preserve
                       compatibility, we compute num_days as len(activities_by_day) - 1.

    Supports up to 8 days by chunking into groups of max 4 days and creating
    multiple agenda slides when needed.
    """

    # Step 1: Get the main color from the logo
    color_thief = ColorThief(inputs.input_0_A)
    main_color = color_thief.get_color(quality=1)

    # Load the agenda templates presentation
    pres = Presentation("liste_AGENDAS.pptx")

    # Backward-compat: previous code expected an extra trailing element
    total_days = int(len(activities_by_day)) - 1
    if total_days < 2:
        # For robustness, do nothing if fewer than 2 days (no agenda slide)
        return None

    # Function to create a single agenda slide for up to 4 days
    def create_agenda_slide(days_slice, day_offset):
        num_days = len(days_slice)
        if num_days == 2:
            template_slide = pres.slides[0]
        elif num_days == 3:
            template_slide = pres.slides[1]
        elif num_days == 4:
            template_slide = pres.slides[2]
        else:
            # For num_days == 1, fall back to 2-day template and leave blanks
            template_slide = pres.slides[0]

        slide_layout = pasteIntoPres.slide_layouts.get_by_name("Blank")
        new_slide = pasteIntoPres.slides.add_slide(slide_layout)

        # Copy all shapes from the selected template slide
        for shape in template_slide.shapes:
            el = shape.element
            newel = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        # Replace the activities in the new slide
        for local_idx, activities in enumerate(days_slice):
            # day_idx visible on the slide starts at 1
            day_idx_on_slide = local_idx + 1
            for shape in new_slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            marker_prefix = f"ACTIVITE_{day_idx_on_slide}_"
                            if marker_prefix in run.text:
                                # Extract the column index
                                try:
                                    col_idx = int(run.text.split("_")[-1]) - 1
                                except Exception:
                                    col_idx = 0
                                run.text = activities[col_idx] if col_idx < len(activities) else ""

        # Replace the logo image (first image on the slide)
        for shape in new_slide.shapes:
            if shape.shape_type == 13:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                shape._element.getparent().remove(shape._element)
                new_slide.shapes.add_picture(inputs.input_0_A, left, top, width, height)
                break

        # Adjust non-black shapes to match brand color
        change_shape_colors(new_slide, main_color)

        return new_slide

    # Chunk into groups of up to 4 days and create multiple agenda slides if needed
    # activities_by_day[:-1] removes the trailing element from legacy usage
    normalized_days = activities_by_day[:-1]
    for chunk_start in range(0, len(normalized_days), 4):
        chunk = normalized_days[chunk_start:chunk_start + 4]
        create_agenda_slide(chunk, chunk_start)

    # Function adds one or more slides directly to pasteIntoPres
    return None